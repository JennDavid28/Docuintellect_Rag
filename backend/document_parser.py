import base64
import logging
import os
import re
import tempfile
import time
from typing import Dict, List, Tuple

logger = logging.getLogger(__name__)

CHUNK_TOKENS = int(os.environ.get("CHUNK_TOKENS", "512"))
CHUNK_OVERLAP = int(os.environ.get("CHUNK_OVERLAP", "128"))
MAX_IMAGE_BYTES = int(os.environ.get("MAX_IMAGE_BYTES", str(4 * 1024 * 1024)))
WHISPER_MODEL = os.environ.get("WHISPER_MODEL", "base")

# ---------------------------------------------------------------------------
# Whisper model — loaded once and cached for the lifetime of the process.
# Loading it on every upload/question would waste 3-10 seconds each time.
# ---------------------------------------------------------------------------
_whisper_model = None

def get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        try:
            import whisper
            logger.info("Loading Whisper model '%s' — this happens once at startup.", WHISPER_MODEL)
            _whisper_model = whisper.load_model(WHISPER_MODEL)
            logger.info("Whisper model loaded successfully.")
        except Exception as exc:
            logger.warning("Could not load Whisper model: %s", exc)
            _whisper_model = None
    return _whisper_model


def estimate_tokens(text: str) -> List[str]:
    """Return list of whitespace-separated words (used for chunking)."""
    return re.findall(r"\S+", text or "")


def chunk_text_with_overlap(text: str, chunk_tokens: int = CHUNK_TOKENS, overlap: int = CHUNK_OVERLAP) -> List[str]:
    words = estimate_tokens(text)
    if not words:
        return []

    chunk_tokens = max(64, chunk_tokens)
    overlap = max(0, min(overlap, chunk_tokens - 1))
    step = chunk_tokens - overlap
    chunks = []

    for start in range(0, len(words), step):
        chunk_words = words[start:start + chunk_tokens]
        if chunk_words:
            chunks.append(" ".join(chunk_words))
        if start + chunk_tokens >= len(words):
            break

    return chunks


def make_page(page_number: int, text: str = "", tables=None, images=None, source: str = "parser") -> Dict:
    tables = tables or []
    images = images or []
    sections = []

    if text.strip():
        sections.append(text.strip())

    for idx, table in enumerate(tables, start=1):
        table_md = table.get("markdown", "").strip()
        if table_md:
            sections.append(f"[Table {idx}]\n{table_md}")

    for idx, image in enumerate(images, start=1):
        description = image.get("description", "").strip()
        if description:
            sections.append(f"[Image {idx} description]\n{description}")

    combined = "\n\n".join(sections)
    return {
        "page_number": page_number,
        "text": combined,
        "tables": tables,
        "images": images,
        "source": source,
    }


def describe_image_with_groq(image_b64: str, mime_type: str = "image/png") -> str:
    api_key = os.environ.get("GROQ_API_KEY", "").strip()
    if not api_key or len(api_key) < 20:
        return "Image extracted from the document. Set GROQ_API_KEY to enable visual description."

    try:
        from groq import Groq

        vision_model = os.environ.get("GROQ_VISION_MODEL", "meta-llama/llama-4-scout-17b-16e-instruct")
        client = Groq(api_key=api_key, timeout=30.0)
        print(f"[Groq Vision] Describing image ({mime_type})...", flush=True)
        _vstart = time.time()
        response = client.chat.completions.create(
            model=vision_model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Describe this document image for retrieval. "
                                "Capture visible text, chart/table structure, labels, numbers, and any facts. "
                                "Be concise but do not omit important details."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{mime_type};base64,{image_b64}"},
                        },
                    ],
                }
            ],
            max_tokens=512,
        )
        print(f"[Groq Vision] Done in {time.time() - _vstart:.1f}s", flush=True)
        return response.choices[0].message.content.strip()
    except Exception as exc:
        logger.warning("Groq Vision image description failed: %s", exc)
        print(f"[Groq Vision] FAILED: {exc}", flush=True)
        return "Image extracted from the document, but visual description failed."


def image_bytes_to_base64(raw: bytes) -> str:
    return base64.b64encode(raw).decode("ascii")


def extract_pdf_images_pypdf(page, max_images: int = 3) -> List[Dict]:
    images = []
    try:
        for idx, image_file in enumerate(getattr(page, "images", [])[:max_images], start=1):
            raw = image_file.data
            if not raw or len(raw) > MAX_IMAGE_BYTES:
                continue
            name = getattr(image_file, "name", f"image_{idx}")
            ext = name.rsplit(".", 1)[-1].lower() if "." in name else "png"
            mime_type = "image/jpeg" if ext in {"jpg", "jpeg"} else "image/png"
            image_b64 = image_bytes_to_base64(raw)
            images.append({
                "name": name,
                "mime_type": mime_type,
                "base64": image_b64,
                "description": describe_image_with_groq(image_b64, mime_type),
            })
    except Exception as exc:
        logger.debug("PDF image extraction skipped: %s", exc)
    return images


def extract_markdown_tables(markdown: str) -> Tuple[str, List[Dict]]:
    lines = markdown.splitlines()
    tables = []
    kept_lines = []
    idx = 0

    while idx < len(lines):
        line = lines[idx]
        next_line = lines[idx + 1] if idx + 1 < len(lines) else ""
        looks_like_table = "|" in line and re.match(r"^\s*\|?[\s:\-|]+\|[\s:\-|]*$", next_line)

        if looks_like_table:
            table_lines = [line, next_line]
            idx += 2
            while idx < len(lines) and "|" in lines[idx].strip():
                table_lines.append(lines[idx])
                idx += 1
            tables.append({"markdown": "\n".join(table_lines).strip()})
            kept_lines.append(f"[Table {len(tables)} extracted separately]")
            continue

        kept_lines.append(line)
        idx += 1

    return "\n".join(kept_lines), tables


def parse_pdf_with_docling(file_path: str) -> List[Dict]:
    # Docling's OCR pipeline (RapidOCR/torch) is not viable in this environment —
    # every attempt fails with "Unsupported configuration: torch.PP-OCRv6.det.small"
    # and falls back to pypdf anyway, but only after paying the cost of initializing
    # the OCR pipeline first. Skip straight to pypdf instead of attempting Docling.
    return []


def parse_pdf(file_path):
    docling_pages = parse_pdf_with_docling(file_path)
    if docling_pages:
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            for idx, page in enumerate(reader.pages):
                images = extract_pdf_images_pypdf(page)
                if images:
                    docling_pages.append(make_page(idx + 1, "", images=images, source="pypdf-images"))
        except Exception as exc:
            logger.debug("PDF image augmentation skipped after Docling parse: %s", exc)
        return docling_pages

    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        pages = []
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            images = extract_pdf_images_pypdf(page)
            pages.append(make_page(idx + 1, text, images=images, source="pypdf"))
        return pages or [make_page(1, "", source="pypdf")]
    except Exception as e:
        return [make_page(1, f"Error parsing PDF: {str(e)}")]

_TOTAL_ROW_PATTERN = re.compile(r'(?i)\btotal\b')

def _split_table_rows(markdown: str, rows_per_chunk: int = 3):
    """Split a table's markdown into small row-groups so each chunk covers only
    a handful of records instead of the entire table. Without this, a 20-row
    table becomes ONE search chunk containing every person/entity in it —
    making it impossible to retrieve or cite a specific row precisely."""
    lines = markdown.splitlines()
    if len(lines) < 3:
        return [markdown]

    header, separator = lines[0], lines[1]
    data_rows = lines[2:]
    if not data_rows:
        return [markdown]

    groups = []
    for i in range(0, len(data_rows), rows_per_chunk):
        chunk_rows = data_rows[i:i + rows_per_chunk]
        groups.append("\n".join([header, separator] + chunk_rows))

    # Aggregate rows ("TOTAL", "COMPANY TOTAL", "Grand Total") are exactly what
    # answers overall-sum questions. Give every total/summary row its OWN dedicated
    # single-row chunk — but ONLY when the total row is not already alone in its
    # group (i.e., it shares a group with other rows). This prevents duplicates.
    total_rows_already_solo = set()
    for i in range(0, len(data_rows), rows_per_chunk):
        chunk_rows = data_rows[i:i + rows_per_chunk]
        if len(chunk_rows) == 1 and _TOTAL_ROW_PATTERN.search(chunk_rows[0]):
            total_rows_already_solo.add(chunk_rows[0])

    for row in data_rows:
        if _TOTAL_ROW_PATTERN.search(row) and row not in total_rows_already_solo:
            groups.append("\n".join([header, separator, row]))

    return groups

def table_to_markdown(rows: List[List[str]]) -> str:
    cleaned = [[str(cell or "").replace("\n", " ").strip() for cell in row] for row in rows if row]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    normalized = [row + [""] * (width - len(row)) for row in cleaned]
    header = normalized[0]
    separator = ["---"] * width
    body = normalized[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def parse_docx(file_path):
    try:
        import docx

        doc = docx.Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        tables = []
        for table in doc.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            markdown = table_to_markdown(rows)
            if markdown:
                tables.append({"markdown": markdown})

        text = "\n".join(paragraphs)
        return [make_page(1, text, tables=tables, source="python-docx")]
    except Exception as e:
        return [make_page(1, f"Error parsing DOCX: {str(e)}")]


def parse_pptx(file_path):
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides = []
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            images = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    raw = shape.image.blob
                    if raw and len(raw) <= MAX_IMAGE_BYTES:
                        mime_type = shape.image.content_type or "image/png"
                        image_b64 = image_bytes_to_base64(raw)
                        images.append({
                            "name": f"slide_{idx + 1}_image",
                            "mime_type": mime_type,
                            "base64": image_b64,
                            "description": describe_image_with_groq(image_b64, mime_type),
                        })
            text = "\n".join(slide_text) or "[Empty Slide]"
            slides.append(make_page(idx + 1, text, images=images, source="python-pptx"))
        return slides
    except Exception as e:
        return [make_page(1, f"Error parsing PPTX: {str(e)}")]


def parse_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [make_page(1, text, source="text")]
    except Exception as e:
        return [make_page(1, f"Error parsing TXT: {str(e)}")]


# ---------------------------------------------------------------------------
# Audio transcription via local Whisper
# Handles: mp4, mp3, wav, webm, ogg, m4a — anything ffmpeg can decode
# ---------------------------------------------------------------------------
def transcribe_with_whisper(file_path: str, filename: str) -> str:
    """
    Run local Whisper on any audio/video file.
    Returns the full transcript text, or an empty string if transcription failed.
    """
    model = get_whisper_model()
    if model is None:
        return ""

    try:
        print(f"[Whisper] Transcribing '{filename}'... (this can take 1-4x the audio length)", flush=True)
        _start = time.time()
        result = model.transcribe(file_path, fp16=False)
        elapsed = time.time() - _start
        segments = result.get("segments") or []
        if segments:
            transcript = "\n".join(
                f"[{float(seg.get('start', 0.0)):.1f}s]: {seg.get('text', '').strip()}"
                for seg in segments
                if seg.get("text", "").strip()
            )
        else:
            transcript = (result.get("text") or "").strip()
        logger.info("Whisper transcription succeeded for '%s' (%d chars, %.1fs).", filename, len(transcript), elapsed)
        print(f"[Whisper] Done in {elapsed:.1f}s — {len(transcript)} chars transcribed.", flush=True)
        return transcript
    except Exception as exc:
        logger.warning("Whisper transcription failed for '%s': %s", filename, exc)
        print(f"[Whisper] FAILED for '{filename}': {exc}", flush=True)
        return ""


def parse_mp4(file_path):
    """Parse any local audio/video file using local Whisper."""
    filename = os.path.basename(file_path)
    transcript = transcribe_with_whisper(file_path, filename)

    if transcript:
        return [make_page(1, f"[Whisper transcript: {filename}]\n\n{transcript}", source="whisper")]

    fallback = (
        f"[Video/Audio Processing Details]\n"
        f"Filename: {filename}\n"
        f"Status: File uploaded, but Whisper transcription is not available.\n\n"
        f"Make sure openai-whisper is installed and ffmpeg is on your system PATH."
    )
    return [make_page(1, fallback, source="mp4")]


# ---------------------------------------------------------------------------
# YouTube ingestion
# Strategy:
#   1. Try youtube-transcript-api first (fast, free, no download needed).
#   2. If captions are unavailable/disabled, download the audio with yt-dlp
#      and transcribe with local Whisper — so NO video is ever missing.
# ---------------------------------------------------------------------------
def get_youtube_id(url: str):
    pattern = (
        r'(?:https?:\/\/)?(?:www\.)?'
        r'(?:youtube\.com\/(?:[^\/\n\s]+\/\S+\/|(?:v|e(?:mbed)?)\/|\S*?[?&]v=)|youtu\.be\/)'
        r'([a-zA-Z0-9_-]{11})'
    )
    match = re.search(pattern, url)
    return match.group(1) if match else None


def _fetch_youtube_captions(video_id: str):
    """Fetch captions with language preference: hi → en → any available.
    Previously the code just grabbed the first track returned by YouTube,
    which could be Urdu, auto-generated Arabic, etc. even for Hindi videos.
    Now we explicitly try preferred languages first so Hindi videos get Hindi
    transcripts and the downstream summarisation/RAG works correctly."""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi

        PREFERRED_LANGS = ["hi", "en", "en-IN", "en-GB"]

        def _items_to_list(fetched):
            try:
                return [{"start": item.start, "text": item.text} for item in fetched]
            except Exception:
                # Older API versions return plain dicts
                return [{"start": item.get("start", 0), "text": item.get("text", "")} for item in fetched]

        try:
            # New API (>=0.6)
            transcript_list = YouTubeTranscriptApi().list(video_id)
            available = list(transcript_list)

            # Build a lookup: language_code → transcript object
            by_lang = {t.language_code: t for t in available}

            # Try preferred languages in order
            for lang in PREFERRED_LANGS:
                if lang in by_lang:
                    logger.info("YouTube: fetching '%s' captions for video_id=%s", lang, video_id)
                    fetched = by_lang[lang].fetch()
                    return _items_to_list(fetched)

            # Fall back to the first available track (whatever language it is)
            if available:
                fallback = available[0]
                logger.info(
                    "YouTube: no preferred language found, falling back to '%s' for video_id=%s",
                    fallback.language_code, video_id,
                )
                fetched = fallback.fetch()
                return _items_to_list(fetched)

        except AttributeError:
            # Old API fallback — try language codes then grab first
            transcripts = YouTubeTranscriptApi.list_transcripts(video_id)
            for lang in PREFERRED_LANGS:
                try:
                    return transcripts.find_transcript([lang]).fetch()
                except Exception:
                    continue
            transcript = next(iter(transcripts))
            return transcript.fetch()

    except Exception as exc:
        logger.info("YouTube captions unavailable for %s: %s", video_id, exc)
        return None


def _download_youtube_audio(url: str) -> str:
    """
    Download the audio track of a YouTube video using yt-dlp.
    Returns path to the downloaded .mp3 temp file, or empty string on failure.
    Requires yt-dlp to be installed: pip install yt-dlp
    """
    try:
        import yt_dlp  # noqa: F401 — just check it's importable
    except ImportError:
        logger.warning(
            "yt-dlp is not installed. Cannot download YouTube audio for Whisper transcription. "
            "Run: pip install yt-dlp"
        )
        return ""

    tmp_dir = tempfile.mkdtemp()
    output_template = os.path.join(tmp_dir, "audio.%(ext)s")

    logger.info("YouTube: no captions found — downloading audio via yt-dlp (this can take a while)...")
    print(f"[YouTube] Downloading audio for transcription... (no captions available)", flush=True)
    _dl_start = time.time()

    ydl_opts = {
        "format": "bestaudio/best",
        "outtmpl": output_template,
        "postprocessors": [
            {
                "key": "FFmpegExtractAudio",
                "preferredcodec": "mp3",
                "preferredquality": "128",
            }
        ],
        "quiet": True,
        "no_warnings": True,
        "socket_timeout": 60,      # increase from default 20s
        "retries": 5,              # retry 5 times before giving up
        "fragment_retries": 5,
        "cookiefile": os.environ.get("YOUTUBE_COOKIES_FILE", "backend/www.youtube.com_cookies.txt"),
    }

    try:
        import yt_dlp

        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])

        audio_path = os.path.join(tmp_dir, "audio.mp3")
        if os.path.exists(audio_path):
            logger.info("yt-dlp audio download succeeded: %s", audio_path)
            print(f"[YouTube] Audio download finished in {time.time() - _dl_start:.1f}s", flush=True)
            return audio_path

        # yt-dlp may have used a different extension
        for fname in os.listdir(tmp_dir):
            full = os.path.join(tmp_dir, fname)
            if os.path.isfile(full):
                return full

        logger.warning("yt-dlp ran but no output file found in %s", tmp_dir)
        return ""
    except Exception as exc:
        logger.warning("yt-dlp download failed for %s: %s", url, exc)
        print(f"[YouTube] Audio download FAILED: {exc}", flush=True)
        return ""


def parse_youtube(url: str) -> List[Dict]:
    """
    Ingest a YouTube video:
      Step 1 — Try existing captions via youtube-transcript-api (instant).
      Step 2 — If no captions, download audio via yt-dlp and transcribe with
                local Whisper (slower but guarantees ingestion).
    """
    video_id = get_youtube_id(url)
    if not video_id:
        return [make_page(
            1,
            f"[YouTube Video: {url}]\nError: Could not extract a valid YouTube video ID from this URL.",
            source="youtube",
        )]

    # ---- Step 1: captions ----
    caption_items = _fetch_youtube_captions(video_id)
    if caption_items:
        text = "\n".join(
            f"[{float(item['start']):.1f}s]: {item['text']}"
            for item in caption_items
        )
        logger.info("YouTube captions fetched for video_id=%s (%d segments).", video_id, len(caption_items))
        return [make_page(
            1,
            f"[YouTube Video Transcript (captions): {url}]\n\n{text}",
            source="youtube-captions",
        )]

    # ---- Step 2: Whisper fallback ----
    logger.info(
        "No captions found for video_id=%s. Attempting yt-dlp + Whisper transcription.", video_id
    )
    audio_path = _download_youtube_audio(url)

    if audio_path:
        try:
            transcript = transcribe_with_whisper(audio_path, f"youtube_{video_id}")
        finally:
            # Clean up the temp audio file regardless of success/failure
            try:
                os.remove(audio_path)
                os.rmdir(os.path.dirname(audio_path))
            except Exception:
                pass

        if transcript:
            return [make_page(
                1,
                f"[YouTube Video Transcript (Whisper): {url}]\n\n{transcript}",
                source="youtube-whisper",
            )]

    # ---- Both methods failed ----
    return [make_page(
        1,
        (
            f"[YouTube Video: {url}]\n\n"
            "Could not extract transcript by either method:\n"
            "  1. youtube-transcript-api: no captions available on this video.\n"
            "  2. yt-dlp + Whisper: either yt-dlp is not installed, ffmpeg is missing, "
            "or Whisper failed to transcribe the audio.\n\n"
            "To fix method 2, run:\n"
            "  pip install yt-dlp\n"
            "  # and make sure ffmpeg is on your PATH"
        ),
        source="youtube",
    )]


def build_search_chunks(pages: List[Dict], chunk_tokens: int = CHUNK_TOKENS, overlap: int = CHUNK_OVERLAP) -> List[Dict]:
    chunks = []
    for page in pages:
        page_number = page.get("page_number", 1)

        for table_idx, table in enumerate(page.get("tables", []), start=1):
            markdown = table.get("markdown", "").strip()
            if markdown:
                row_groups = _split_table_rows(markdown, rows_per_chunk=3)
                for group_idx, group_markdown in enumerate(row_groups, start=1):
                    chunks.append({
                        "page_number": page_number,
                        "content_type": "table",
                        "text_content": f"[Table {table_idx} on page {page_number}, rows group {group_idx}]\n{group_markdown}",
                        "table_markdown": group_markdown,
                        "image_count": 0,
                    })

        for image_idx, image in enumerate(page.get("images", []), start=1):
            description = image.get("description", "").strip()
            if description:
                chunks.append({
                    "page_number": page_number,
                    "content_type": "image",
                    "text_content": f"[Image {image_idx} on page {page_number}]\n{description}",
                    "image_description": description,
                    "image_count": 1,
                })

        plain_text = page.get("text", "").strip()
        for chunk in chunk_text_with_overlap(plain_text, chunk_tokens, overlap):
            chunks.append({
                "page_number": page_number,
                "content_type": "text",
                "text_content": chunk,
                "image_count": len(page.get("images", [])),
            })

    return chunks


def extract_document_pages(file_path: str, file_type: str) -> List[Dict]:
    file_type = file_type.lower().strip()
    if file_type == "pdf":
        return parse_pdf(file_path)
    elif file_type in ["docx", "doc"]:
        return parse_docx(file_path)
    elif file_type in ["pptx", "ppt"]:
        return parse_pptx(file_path)
    elif file_type in ["txt", "md"]:
        return parse_txt(file_path)
    elif file_type in ["mp4", "mp3", "wav", "webm", "ogg", "m4a"]:
        # All audio/video formats go through local Whisper
        return parse_mp4(file_path)
    elif file_type == "youtube":
        return parse_youtube(file_path)
    else:
        return parse_txt(file_path)

def get_youtube_title(url: str) -> str:
    """Fetch the actual video title for auto-naming uploaded YouTube links.
    Tries YouTube's public oEmbed endpoint first (fast, no auth, not subject
    to the same bot-detection as scraping captions/audio). Falls back to
    yt-dlp metadata-only extraction (no download) if oEmbed fails, and to the
    raw URL as a last resort so upload never breaks because of this."""
    try:
        import requests
        resp = requests.get(
            "https://www.youtube.com/oembed",
            params={"url": url, "format": "json"},
            timeout=8,
        )
        if resp.ok:
            title = (resp.json() or {}).get("title", "").strip()
            if title:
                return title
    except Exception as exc:
        logger.info("oEmbed title fetch failed for %s: %s", url, exc)

    try:
        import yt_dlp
        ydl_opts = {"quiet": True, "no_warnings": True, "skip_download": True}
        cookiefile = os.environ.get("YOUTUBE_COOKIES_FILE", "backend/youtube_cookies.txt")
        if os.path.exists(cookiefile):
            ydl_opts["cookiefile"] = cookiefile
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            title = (info or {}).get("title", "").strip()
            if title:
                return title
    except Exception as exc:
        logger.info("yt-dlp title fetch failed for %s: %s", url, exc)

    return url